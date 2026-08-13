#!/usr/bin/env python3
"""build_deck.py - render the whitepaper figures to PNG and assemble a leadership
PowerPoint that follows the paper's argumentation, with speaking narrative in the
speaker notes. Run with the clean_env python (has python-pptx)."""
from __future__ import annotations

import os
import re
import subprocess
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

PAPER = "/home/jlmorg1/work/Johnny_AI_ML_LLM_experiments_2024/papers/ai-metadata-session-data-repository"
OUT_DIR = "/home/jlmorg1/.openclaw/workspace/deck_build"
IMG_DIR = os.path.join(OUT_DIR, "figs")
PPTX = os.path.join(OUT_DIR, "AI_Metadata_Repository_Briefing.pptx")

os.makedirs(IMG_DIR, exist_ok=True)

NAVY = RGBColor(0x11, 0x2A, 0x46)
ACCENT = RGBColor(0x1F, 0x6F, 0x8B)
INK = RGBColor(0x20, 0x24, 0x28)
MUT = RGBColor(0x5A, 0x63, 0x6B)
LIGHT = RGBColor(0xF3, 0xF5, 0xF7)

STANDALONE = r"""\documentclass[border=8pt]{standalone}
\usepackage{times}
\usepackage{tikz}
\usepackage{pgfplots}\pgfplotsset{compat=1.18}
\usetikzlibrary{arrows.meta, positioning, shapes.geometric, fit, backgrounds, matrix, calc, decorations.pathreplacing}
\usepackage{amsmath}
\renewcommand{\ref}[1]{}
\begin{document}
\input{FIG}
\end{document}
"""


def render_figures():
    """Compile each figure standalone and rasterize to PNG. Returns {name: pngpath}."""
    figs = {
        "fig1": "fig1_architecture_flow.tex",
        "fig2": "fig2_data_model.tex",
        "fig3": "fig3_measurement.tex",
        "fig4": "fig4_lifecycle.tex",
        "fig5": "fig5_deployment.tex",
        "fig6": "fig6_forensics.tex",
    }
    out = {}
    for name, fn in figs.items():
        src = os.path.join(PAPER, fn)
        with open(src, encoding="utf-8") as f:
            content = f.read()
        # section refs are noise in a standalone figure; drop the parenthetical
        content = re.sub(r"\s*\(Section~\\ref\{[^}]*\}\)", "", content)
        content = re.sub(r"Section~\\ref\{[^}]*\}", "the linked record", content)
        with tempfile.TemporaryDirectory() as td:
            with open(os.path.join(td, f"{name}.tex"), "w", encoding="utf-8") as f:
                f.write(content)
            with open(os.path.join(td, "wrap.tex"), "w", encoding="utf-8") as f:
                f.write(STANDALONE.replace("FIG", name))
            r = subprocess.run(
                ["pdflatex", "-interaction=nonstopmode", "-halt-on-error", "wrap.tex"],
                cwd=td, capture_output=True, text=True)
            pdf = os.path.join(td, "wrap.pdf")
            if not os.path.exists(pdf):
                print(f"  [FAIL] {name}\n{r.stdout[-1500:]}")
                continue
            png = os.path.join(IMG_DIR, f"{name}.png")
            subprocess.run(["pdftoppm", "-png", "-r", "300", "-singlefile", pdf,
                            os.path.join(IMG_DIR, name)], check=True)
            out[name] = png
            print(f"  [ok] {name} -> {png}")
    return out


def img_aspect(path):
    try:
        from PIL import Image
        with Image.open(path) as im:
            return im.width / im.height
    except Exception:
        return 2.5  # wide-diagram default


# ---------------------------------------------------------------- slide content
# Each slide: title, bullets (list of (text, level)), image key or None, notes.
def build(figs):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_title(slide, text, sub=None):
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = text
        r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = NAVY
        r.font.name = "Calibri"
        if sub:
            p2 = tf.add_paragraph()
            r2 = p2.add_run(); r2.text = sub
            r2.font.size = Pt(15); r2.font.color.rgb = ACCENT; r2.font.name = "Calibri"
        # accent rule
        ln = slide.shapes.add_textbox(Inches(0.62), Inches(1.28), Inches(3.2), Inches(0.05))
        return tb

    def add_bullets(slide, bullets, top, height, width=Inches(12.1), left=Inches(0.7), size=18):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame; tf.word_wrap = True
        first = True
        for text, lvl in bullets:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = lvl
            p.space_after = Pt(6)
            bullet = ("•  " if lvl == 0 else "–  ")
            r = p.add_run(); r.text = bullet + text
            r.font.size = Pt(size - lvl * 2)
            r.font.color.rgb = INK if lvl == 0 else MUT
            r.font.name = "Calibri"
            if lvl == 0:
                r.font.bold = False
        return tb

    def add_image(slide, key, top, maxh):
        if not key or key not in figs:
            return
        path = figs[key]
        asp = img_aspect(path)
        maxw = Inches(12.2)
        w = maxw
        h = Inches(w.inches / asp)
        if h > maxh:
            h = maxh
            w = Inches(h.inches * asp)
        left = Inches((13.333 - w.inches) / 2)
        slide.shapes.add_picture(path, left, top, width=w, height=h)

    def add_notes(slide, text):
        slide.notes_slide.notes_text_frame.text = text.strip()

    def slide(title, bullets=None, image=None, notes="", sub=None, title_only=False):
        s = prs.slides.add_slide(blank)
        # subtle side band
        band = s.shapes.add_textbox(Inches(0), Inches(0), Inches(0.22), Inches(7.5))
        band.fill  # noqa
        add_title(s, title, sub)
        if image:
            if bullets:
                add_bullets(s, bullets, Inches(1.5), Inches(1.7), size=17)
                add_image(s, image, Inches(3.35), Inches(3.75))
            else:
                add_image(s, image, Inches(1.7), Inches(5.3))
        elif bullets:
            add_bullets(s, bullets, Inches(1.7), Inches(5.3))
        add_notes(s, notes)
        return s

    B = lambda t, l=0: (t, l)

    # 1 TITLE ---------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.7), Inches(2.6))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run()
    r.text = "AI Metadata and Session Data Repository"
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = NAVY; r.font.name = "Calibri"
    p2 = tf.add_paragraph(); r2 = p2.add_run()
    r2.text = "A Client Architecture for Managing, Measuring, and Mitigating Enterprise Generative AI"
    r2.font.size = Pt(20); r2.font.color.rgb = ACCENT; r2.font.name = "Calibri"
    p3 = tf.add_paragraph(); p3.space_before = Pt(18); r3 = p3.add_run()
    r3.text = "Johnny Morgan  |  UMBC Department of Information Systems"
    r3.font.size = Pt(15); r3.font.color.rgb = MUT; r3.font.name = "Calibri"
    add_notes(s, """
Thank you. The argument I want to leave you with is one sentence: you cannot manage what you cannot see, and seeing how this organization uses generative AI requires capturing that use before we know what we will need to ask of it. Everything else follows from that. Over the next few minutes I will show the problem we are exposed to today, why the fix has to be a record and not a dashboard, the architecture that holds that record, and what it unlocks: measurable return, detectable and mitigatable harm, and a defensible account when something goes wrong. This is a proposal and a theory of operation. I am not reporting results; I am asking us to build the substrate that makes results measurable.
""")

    # 2 PROBLEM -------------------------------------------------------------
    slide("You cannot manage AI you cannot see",
          [B("A model now drafts customer replies, writes code, summarizes documents, and helps decide"),
           B("Every use carries a promise and a hazard: fluent, confident output that can be wrong"),
           B("The output is used; the conditions that produced it are gone"),
           B("An uninstrumented deployment cannot be measured, governed, or defended")],
          notes="""
Adoption is running ahead of management. A model is already drafting contract clauses, answering customers, writing code, and summarizing case files across the organization. Each of those interactions leaves almost nothing behind: we keep the output and lose the conditions that produced it. That is the gap. It is not a missing dashboard, it is the absence of a record. And without the record two things are invisible at once: the value we are getting, which we cannot prove, and the failures we are exposed to, which we cannot see. Governance frameworks and regulation have started to treat the record of how AI is used as an obligation, not an option. Most deployments, including most of ours, keep no such record.
""")

    # 3 SIX QUESTIONS -------------------------------------------------------
    slide("Six questions leadership cannot answer today",
          [B("Inventory: which AI tools, models, and versions are in use, by whom"),
           B("Utilization and cost: sessions, tokens, and spend, by team and workflow"),
           B("Effectiveness: task success, rework, acceptance, escalation"),
           B("Risk: how often the model fabricates, whether it recurs, whether it cascades"),
           B("In-session degradation: context rot as a session lengthens"),
           B("The people and the model: workforce skill, and the model's own drift over time")],
          notes="""
Here is the concrete test. A manager accountable for an AI deployment would want to establish each of these, and today usually cannot. What do we even have deployed, and at what version. How much are we using it and what does it cost. Is it actually doing useful work. How often does it produce something wrong, and does that wrong thing spread. Does quality decay as a session gets long. And how skillfully is our workforce using it, while the model itself shifts underneath us. The instinct is to answer these by adding metrics. That instinct fails, for a reason I will give on the next slide: the failures that matter most are the ones nobody thought to declare in advance.
""")

    # 4 LIABILITY HOOK ------------------------------------------------------
    slide("Why the record is indispensable",
          [B("AI harms are foreseeable: prompt efficacy is known not to be guaranteed"),
           B("They are detectable: the interaction that carried the fault can be examined"),
           B("They are mitigatable: a detected fault can be corrected before it spreads"),
           B("A harm you could have foreseen, detected, and mitigated, but did not record, is one you own"),
           B("The architecture turns a realized loss into a managed event", 1)],
          notes="""
This is the slide that matters most to leadership, because it is about liability. An AI harm is not an act of god. It is foreseeable, because we know a prompt that works today can fail under a different input or a longer context. It is detectable, because the interaction that carried the fault can be examined. And it is mitigatable, because a fault we catch can be corrected before it spreads. Now put those together. An organization that could have foreseen, detected, and mitigated a harm, and did not, because it kept no record, owns the consequence. The record is what converts a foreseeable, detectable, mitigatable harm from a realized loss into a managed one. That is the difference between explaining an incident and being blindsided by it.
""")

    # 5 THREE CASES ---------------------------------------------------------
    slide("Three cases that make it concrete",
          [B("A prompt that quietly stops working"),
           B("A customer's phrasing pushes the model outside its system prompt; it states a policy we do not hold. With the record, the deviation is caught before it recurs.", 1),
           B("A fabrication that enters the record as fact"),
           B("A model adds a figure the source does not contain; it moves into a report, then a decision. With provenance, it is traceable and the correction reaches everywhere it went.", 1),
           B("A cascade across a long session"),
           B("The model forgets an earlier constraint; later outputs violate it while each looks locally fine. With session capture, the drift is flagged before the output is trusted.", 1)],
          notes="""
Three short scenarios, each a real failure shape. First, a system prompt that silently stops constraining the model when a customer phrases things a certain way, and the assistant states a policy we do not hold. Without the record, we learn when a customer acts on it. Second, the sharpest one: a fabricated figure that is indistinguishable from a sourced one, that moves into a report and then a decision. Provenance is the defense: a generated assertion stays distinguishable from an established fact. Third, a long session where the model forgets an earlier constraint and every later output looks locally reasonable and passes review on its own. In each case the difference between a managed event and a corporate loss is one thing: the record.
""")

    # 6 EPISTEMOLOGY --------------------------------------------------------
    slide("Capture before the question",
          [B("Visibility is a precondition, not a feature: you cannot measure what you did not record"),
           B("The behaviors worth measuring are not known when the data is captured"),
           B("A new failure mode becomes legible only after instances accumulate and someone builds the lens"),
           B("Omitted capture is unrecoverable: a gap is a permanent blind spot, not a delay"),
           B("Capturing now is cheap; reconstructing later is impossible", 1)],
          notes="""
This is the intellectual core, and it is about knowledge, not technology. Two claims. First, visibility is a precondition. A KPI, a measure of effectiveness, a hallucination rate, these are all functions computed over captured data. Where there is no capture, the function has no input and the question cannot be asked at all. Second, the behaviors worth measuring are not known at capture time. A new failure mode only becomes measurable after enough instances accumulate to form a pattern and after someone builds the lens to see it, and both of those happen after capture. Put the two together and you get an asymmetry: capture you omitted before a phenomenon emerged cannot be recovered afterward. There is no retrospective study of data nobody wrote down. So the rational policy is to capture earlier, wider, and at higher fidelity than any present question requires.
""")

    # 7 MONITORING VS OBSERVABILITY ----------------------------------------
    slide("Monitoring answers known questions; observability answers unknown ones",
          [B("Monitoring: declared metrics, thresholds, dashboards, chosen in advance"),
           B("Most current AI tooling is monitoring; it answers questions its designers already hold"),
           B("Observability: record behavior at enough fidelity to answer questions not yet formed"),
           B("Because AI failure modes emerge, the enterprise needs an observability substrate, not a dashboard")],
          notes="""
The distinction that makes this precise is monitoring versus observability. Monitoring instruments a system against questions its designers already hold: fixed metrics, thresholds, dashboards. It answers known unknowns, and it is useful. Most AI operational tooling on the market today is monitoring: it traces calls and scores them against evaluations chosen in advance. Observability is different: it records behavior at enough fidelity that a question nobody anticipated at instrumentation time can still be answered later by interrogating the record. Because AI failure modes emerge, the questions that will matter most are exactly the ones nobody thought to declare. So what an enterprise needs is an observability substrate, and every design choice that follows, immutability, granularity, provenance, retention, is justified by whether it keeps a future, unspecified question answerable.
""")

    # 8 ARCHITECTURE fig1 ---------------------------------------------------
    slide("Reference architecture: a simple flow",
          [B("Instrument the AI systems in use; store their session content"),
           B("Derive metadata into stores separated by sensitivity; draw outcomes and risks; report")],
          image="fig1",
          notes="""
The architecture is simple at the top; its complexity is a drill-down. At the top it is a flow. The AI systems in use are instrumented. Their session content is stored. Metadata is derived from that content into a small set of repositories separated by sensitivity, performance in one, hallucination and risk in another, session continuity in a third. Outcomes and risks are drawn from that metadata, and the whole thing is reported to the organization, attributable by organization, product line, and purpose. The lower layers preserve and organize the record; the upper layers interrogate and report it. Keep that flow in mind, because the rest of the talk drills into each box.
""")

    # 9 DATA MODEL fig2 -----------------------------------------------------
    slide("The captured event: immutable and linked",
          [B("Each interaction is one immutable event: the text and the conditions that produced it"),
           B("System prompt, decoding parameters, model version, tool calls, human edits and acceptances"),
           B("Antecedent and consequent links; a raw append-only form and a queryable structured form")],
          image="fig2",
          notes="""
If an event's meaning is knowable only through its cause and its trajectory, then the unit of capture cannot be a flat log line. Each interaction is stored as one immutable event that records both the exchanged text and the conditions that produced it: the system prompt, the decoding parameters, the model version, the tool calls and results, and the human interactions, the edits, acceptances, and rejections. A hallucination produced at a high temperature under a permissive prompt is a different event, with a different cause, from the same words produced under a constraining one. Events are immutable and linked to their antecedents and consequents, so a cascade can be traversed rather than guessed. And each record is kept twice: a raw, append-only form that is never rewritten, and a structured, queryable form. A later re-classification adds a label; it never alters the original.
""")

    # 10 MEASUREMENT fig3 ---------------------------------------------------
    slide("Measurement: objectives are exogenous and layered",
          [B("The organization defines its measures, in layers: workflow, division, corporate"),
           B("The architecture computes them over the record and rolls outcomes up the hierarchy"),
           B("MOPs, MOEs, KPIs, plus workforce and model acumen and the trust-calibration gap")],
          image="fig3",
          notes="""
A crucial design decision: measures of performance, effectiveness, and KPIs are not properties of the architecture. They are the objectives of the organization, defined in layers, corporate strategy at the top, division and then workflow below. Each layer names its own measures. The architecture does not impose a fixed metric set, and it should not, because objectives change. Its job is to collect, detect, catalog, and store the outcomes those measures are computed from, and to roll them up from the workflow where they occur to the strategic objective they serve. On top of the standard families it adds measures a fixed dashboard would miss: how skillfully the workforce directs and checks the AI, the model's reliability and drift, and the gap between how much people trust the model and how reliable it actually is.
""")

    # 11 PORTFOLIO ----------------------------------------------------------
    slide("Proving the return on the AI investment",
          [B("Cost is known from the record; outcome value from the outcome repository"),
           B("Return is the relation between them, reported at each layer of the objective hierarchy"),
           B("'How far does AI multiply human capital' becomes a computed quantity, not an assertion"),
           B("The record also shows the downside: over-reliance that degrades workforce competency", 1)],
          notes="""
The architecture is a funded investment, and like any portfolio investment it has to justify its cost. Because the record makes both cost and outcome measurable, return becomes calculable rather than asserted. Cost we know from the event record. Outcome value we know from the outcome repository, and because each outcome links to the interactions that produced it, we can attribute value to the work that earned it. The question an executive most wants answered, how far does AI multiply what our workforce produces, becomes a computed quantity: output with AI set against output without, tracked over time. And it is neutral to sign. The same record shows where over-reliance is quietly degrading skills our workforce would otherwise keep, which is a cost a naive productivity number hides.
""")

    # 12 LIFECYCLE fig4 -----------------------------------------------------
    slide("The hallucination lifecycle, from adverse-event surveillance",
          [B("Instance, to tracked incidence, to cause and cascade, to mitigation, to prevention"),
           B("Mitigation efficacy is measured as a change in incidence across a dated boundary"),
           B("An effective mitigation moved upstream becomes prevention")],
          image="fig4",
          notes="""
For harm we borrow a lifecycle from public-health adverse-event surveillance, because the shape is the same. A single failure is an instance, and its importance is not evident at the moment it occurs, because impact is a function of trajectory. Repeated instances form an incidence we can track. Some trace to a cause; some cascade downstream. A mitigation applied after detection has an efficacy that is itself measurable, as a change in incidence across the dated boundary where we applied it. And a mitigation that works and gets moved upstream becomes prevention. Every step here is a query over the linked record, not an inference from a point-in-time log. This is how we turn hallucination from an anecdote into something managed on a curve.
""")

    # 13 TAXONOMY -----------------------------------------------------------
    slide("A working taxonomy of hallucination",
          [B("Fabrication  |  Overconfidence and misrepresentation  |  Staleness"),
           B("Context loss  |  Attribution error  |  Scope and instruction divergence"),
           B("Agentic and tool-use failure"),
           B("Adjacent, and also tracked: refusal and over-restriction (a withheld benefit), and its inverse"),
           B("Fabrication is the sharpest: fiction in the form of fact. Provenance keeps it distinguishable.", 1)],
          notes="""
Hallucination is not one thing, and detecting it requires knowing which form is in play. The taxonomy has fabrication, overconfidence and misrepresentation, staleness, context loss, attribution error, scope and instruction divergence, and agentic or tool-use failure. Each category is detectable only from particular captured fields, which is precisely why the data model records the generating conditions and the provenance and not just the output. Alongside these, though they are not hallucinations, we also track refusal and over-restriction, where the model declines legitimate work and the harm is a withheld benefit, and the inverse, where it complies with something it should have declined. The most damaging form is fabrication: fiction presented in the form of established fact. The architectural defense is provenance. The repository does not stop the model from fabricating; it stops a fabrication from silently becoming part of our fact base.
""")

    # 14 FORENSICS fig6 -----------------------------------------------------
    slide("Harm forensics: the post-incident investigation",
          [B("When a harm occurs, impact and scope are not optional; they are a required investigation"),
           B("Reconstruct the trigger and conditions, attribute the actors, trace the cascade, prescribe a guardrail")],
          image="fig6",
          notes="""
This is the risk surface leadership feels most directly, and it is where the whole architecture pays off at once. When an AI-induced harm occurs, the harm is not the end of the event. It opens an investigation we are obliged to conduct, and the questions are forensic: what was produced, what triggered it, what caused it, how far did it spread, who and what is accountable, and how do we prevent a recurrence. None of these can be answered after the fact for an interaction nobody captured. The record is the forensic substrate. The loop on this slide is the workflow: reconstruct the trigger and conditions, attribute the event, trace its scope and cascade, prescribe a guardrail matched to the cause, and measure whether the guardrail worked. Let me take the two hardest parts, attribution and guardrails, on the next slides.
""")

    # 15 ATTRIBUTION --------------------------------------------------------
    slide("Attribution: insider misuse vs rogue behavior",
          [B("The record separates human turns and human interactions from model turns and tool calls"),
           B("AI actor is the source: the model fabricates, misrepresents, or takes an unsafe autonomous action (rogue / agentic)"),
           B("Human actor is the source: a person directs the AI against the organization (insider misuse; the AI is the instrument)"),
           B("Only the record tells them apart: human intent is legible only through the sequence of instructions and uses")],
          notes="""
A forensic account has to attribute the harm, and attribution starts by separating what the human did from what the AI did. The record supports that because it captures both sides as distinct, linked events. From that we can tell two cases apart that we must never conflate. In the first, the AI actor is the source: the model fabricated, misrepresented, or took an unsafe autonomous action outside its authorized scope. That is the rogue or agentic case. In the second, the human actor is the source: a person directed the AI to act against the organization, and the AI was merely the instrument. That is insider misuse, and its signature is not a model error at all, it is a pattern of instruction and use, extraction at volume, probing of scope or permission. The same record grounds both findings, and only the record lets us distinguish them, because a human actor's intent is legible only through the sequence of instructions the record preserves.
""")

    # 16 GUARDRAILS ---------------------------------------------------------
    slide("From post-mortem to guardrail",
          [B("System-prompt engineering: constrain the instruction frame"),
           B("Input and output validation: check inputs, screen outputs before use"),
           B("Human-in-the-loop review gate: require review before an output is used or an action taken"),
           B("Tool and permission constraint: limit unauthorized autonomous action"),
           B("Workforce hallucination-awareness training: raise the ability to recognize and check failure"),
           B("Model version and configuration control"),
           B("Each guardrail is dated and its efficacy measured, so a nominal control is distinguishable from a working one", 1)],
          notes="""
An investigation is not done when it explains a harm. Its purpose is to prevent the next one, and its output is a guardrail chosen to fit the cause it found. If the cause was a permissive system prompt, the guardrail is prompt engineering. If it was an unchecked input or unscreened output, it is validation. If a human review would have caught it, it is a human-in-the-loop gate, and the record also shows when a gate existed but was not enforced. If an agent took an unsafe action, it is a tool or permission constraint. If the cause was over-trust or an unrecognized fabrication, the guardrail is workforce training. And if it was a model change, it is version control. The key point for accountability: each guardrail is itself recorded and dated, and its efficacy is measured by the same machinery that detected the harm. A guardrail whose efficacy the record cannot demonstrate is, for accountability purposes, indistinguishable from no guardrail at all.
""")

    # 17 OBLIGATIONS + TRUST -----------------------------------------------
    slide("Obligations and calibrated trust",
          [B("An output can be factually right and still violate a policy, a law, a standard, or a procedure"),
           B("Obligation classes: organizational, legal and regulatory, ethical, procedural, professional, factual, safety"),
           B("Trust must be calibrated: earned (measured over the record) and curated (surfaced to the user)"),
           B("Over-trust turns every fabrication into an accepted fact; the record makes over-reliance visible", 1)],
          notes="""
Two connected points. First, correctness is not the only bar. An AI output can be factually right and still violate an internal policy, a law, a professional standard, or a required procedure. So we evaluate AI use against each class of obligation that applies, and that is only possible over a record that holds the output together with its conditions and its downstream use. A violation is just one of the negative outcomes we already track. Second, trust. People tend to treat a fluent, confident model as an expert and extend it trust it has not earned. The goal is not maximal trust or minimal trust, it is calibrated trust: reliance that tracks demonstrated reliability. That requires two things the record provides. Trust has to be earned, which means measured, and it has to be curated, which means the provenance and the model's track record are surfaced to the user at the moment they rely on it.
""")

    # 18 DEPLOYMENT fig5 ----------------------------------------------------
    slide("Deployment, and a future leg",
          [B("Instrument at the AI-use surfaces: an API gateway or proxy, agent hooks, chat and IDE connectors"),
           B("Converge into one governed store, or isolate systems where separation must be guaranteed"),
           B("Future: controlled continuity reinjection and per-workflow lane isolation")],
          image="fig5",
          notes="""
Operationally it is a client-side system that captures at the points of AI use, an API gateway or proxy, hooks in agent frameworks, connectors to chat and development tools, and serves reporting back. It monitors its own health, because an observability system has to be observable; an outage in capture should be seen, not mistaken for quiet. On topology there is a choice: converge the channels into one store with logical separation where that is enough, or isolate systems where a separation must be guaranteed, for instance keeping live operations apart from drills, or fencing a sensitive product line. And there is a deliberate future leg: controlled continuity, resuming a session from a curated summary drawn from the record, and workflow-lane isolation, so a single user's unrelated tasks never bleed into one another, because uncontrolled context mixing is itself a source of error in these systems.
""")

    # 19 GOVERNANCE ---------------------------------------------------------
    slide("Governance, privacy, and compliance",
          [B("Segregation by sensitivity is the primary control, not a convenience; it enforces need-to-know"),
           B("Pseudonymous by default: measure roles and workflows, not people; re-identification is a governed, recorded exception"),
           B("Insider-threat monitoring is bounded by purpose, by access, and by proportionality"),
           B("Supplies the record that NIST AI RMF, ISO 42001 and 23894, the EU AI Act, and OWASP LLM Top 10 presume")],
          notes="""
The record is powerful because it is complete, and sensitive for the same reason, so governance is built into the design rather than bolted on. Segregating the stores by sensitivity is the first control: the person who reports on performance does not thereby gain the organization's risk and insider-threat records. Identity is pseudonymous by default, so we measure roles and workflows, not named people, and re-identifying someone is a governed, recorded exception for a stated purpose such as a security case. The same record is security telemetry, so insider-threat monitoring is legitimate, but it is bounded by purpose, by access, and by proportionality, because surveillance changes the behavior it observes. And crucially for compliance, this architecture supplies exactly the record that the NIST AI Risk Management Framework, ISO 42001 and 23894, the EU AI Act's logging obligations, and the OWASP LLM Top 10 all presume but do not themselves provide.
""")

    # 20 PREDICTIONS --------------------------------------------------------
    slide("Registered predictions: falsifiable, on the record",
          [B("Hallucination recurrence clusters around context-window saturation and compaction, not at random"),
           B("Newly detectable failure classes rise with the age and size of a broadly captured corpus"),
           B("Mitigation efficacy is measurable as an incidence change across a dated boundary"),
           B("Without the record, an incident's impact is systematically underestimated"),
           B("Guardrail efficacy is larger for controls placed upstream of the cause than at the point of harm"),
           B("Over-refusal and under-refusal trade off against each other as the instruction frame tightens")],
          notes="""
Because this is a proposal, I am putting its claims on the record now as dated, falsifiable predictions, so that a deployed system later confirms or refutes them rather than us rationalizing after the fact. A few of them: hallucination recurrence should cluster in time around context-window saturation and compaction, not arrive at random. In a broadly captured corpus, the rate at which we discover new failure classes should rise with the corpus's age and size, precisely because detection is retrospective. Mitigation efficacy should be visible as an incidence change across the dated boundary where we applied it. Without a repository, an incident's true impact should be systematically underestimated, because the cascade is invisible at the moment it occurs. Guardrails placed upstream of a cause should outperform guardrails at the point of harm. And tightening a system prompt should trade over-refusal against under-refusal rather than improve both. These are testable, and that is the point.
""")

    # 21 CLOSE --------------------------------------------------------------
    slide("The ask: build the record first",
          [B("On the record, the harms of AI become foreseeable, detectable, and mitigatable"),
           B("On the same record, AI's value, including how far it multiplies human capital, becomes calculable"),
           B("An organization that keeps it can defend against AI liability and prove its return; one that does not can do neither"),
           B("Capture now is cheap; reconstruction later is impossible. The decision is time-sensitive.", 1)],
          notes="""
To close, back to the one sentence. An organization cannot manage what it cannot see, and it cannot see its use of AI without a record of that use. Build that record and two things become true at once: the harms of AI become foreseeable, detectable, and mitigatable, and the value of AI, including how far it multiplies our people's output, becomes calculable. The same record serves both ends, which means an organization that keeps it can both defend against the liability of AI and prove its return, and an organization that does not can do neither. The one property of this decision I want to leave you with is that it is time-sensitive. Capture we do not start today is capture we can never recover, because there is no retrospective study of data nobody wrote down. The ask is simple: build the record first.
""")

    prs.save(PPTX)
    print(f"\n[deck] {len(prs.slides.__iter__.__self__._sldIdLst)} slides -> {PPTX}")


if __name__ == "__main__":
    print("Rendering figures...")
    figs = render_figures()
    print("Building deck...")
    build(figs)
