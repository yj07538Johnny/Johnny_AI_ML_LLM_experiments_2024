#!/usr/bin/env python3
"""Build the information-need classification PowerPoint deck.

Consumes the PNGs written by make_figures.py and emits a 16:9 .pptx with
speaker notes on every slide.

Usage:
    python make_deck.py [--figures DIR] [--out FILE]
"""

import argparse
import sys
from pathlib import Path

# figkit lives one level up, in papers/, shared across papers in this repo.
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation                                # noqa: E402
from pptx.dml.color import RGBColor                          # noqa: E402
from pptx.enum.shapes import MSO_SHAPE                       # noqa: E402
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN              # noqa: E402
from pptx.util import Inches, Pt                             # noqa: E402

import figkit                                                # noqa: E402

# ---------------------------------------------------------------- style ------
# Hues come from figkit so a colour means the same thing in the deck as in the
# figures: blue is the report side, amber the need side, everywhere.


def C(hexstr) -> RGBColor:
    """figkit hex -> pptx RGBColor."""
    return RGBColor(*figkit.rgb(hexstr))


INK = C(figkit.INK)
MUTED = C(figkit.MUTED)
BLUE = C(figkit.BLUE)
AMBER = C(figkit.AMBER)
GREEN = C(figkit.GREEN)
RED = C(figkit.RED)
VIOLET = C(figkit.VIOLET)
WHITE = C(figkit.PAPER)
SLATE = C(figkit.SLATE_L)

NAVY = C(figkit.INK_DARK)
GREEN_S = C(figkit.GREEN_S)
RED_S = C(figkit.RED_S)
VIOLET_S = C(figkit.VIOLET_S)

BLUE_ON_DARK = C(figkit.BLUE_ON_DARK)
GREEN_ON_DARK = C(figkit.GREEN_ON_DARK)
AMBER_ON_DARK = C(figkit.AMBER_ON_DARK)
VIOLET_ON_DARK = C(figkit.VIOLET_ON_DARK)
RED_ON_DARK = C(figkit.RED_ON_DARK)

AMBER_DEEP = C(figkit.AMBER_DEEP)
VIOLET_DEEP = C(figkit.VIOLET_DEEP)
RED_DEEP = C(figkit.RED_DEEP)
GREEN_DEEP = C(figkit.GREEN_DEEP)

FONT = "Calibri"
W, H = 13.333, 7.5          # inches, 16:9

SECTION_COLORS = {
    "setup": BLUE,
    "supervised": GREEN,
    "pairs": AMBER,
    "attention": VIOLET,
    "cross": RED,
    "train": GREEN,
    "ship": BLUE,
}


# ---------------------------------------------------------------- helpers ----

def new_deck():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, l, t, w, h, color, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t),
                                Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """runs = [(string, size_pt, bold, RGBColor, space_after_pt), ...]"""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (s, sz, bold, col, after) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = s
        p.alignment = align
        p.space_after = Pt(after)
        f = p.font
        f.name = FONT
        f.size = Pt(sz)
        f.bold = bold
        f.color.rgb = col
    return tb


def header(slide, title, sub=None, accent=BLUE, kicker=None):
    rect(slide, 0, 0, 0.22, H, accent)
    y = 0.42
    if kicker:
        text(slide, 0.62, y, 11.8, 0.34,
             [(kicker.upper(), 11, True, accent, 0)])
        y += 0.38
    runs = [(title, 30, True, INK, 4)]
    if sub:
        runs.append((sub, 15, False, MUTED, 0))
    text(slide, 0.62, y, 11.9, 1.25, runs)


def footer(slide, n, label):
    text(slide, 0.62, H - 0.52, 9.0, 0.3, [(label, 10, False, MUTED, 0)])
    text(slide, 11.4, H - 0.52, 1.3, 0.3, [(str(n), 10, False, MUTED, 0)],
         align=PP_ALIGN.RIGHT)


def figure(slide, path, top, max_w=11.4, max_h=4.5, left=None):
    """Insert scaled to fit inside (max_w, max_h). Centred unless `left` given."""
    pic = slide.shapes.add_picture(str(path), Inches(0), Inches(top),
                                   width=Inches(max_w))
    if pic.height > Inches(max_h):
        scale = Inches(max_h) / pic.height
        pic.height = Inches(max_h)
        pic.width = int(pic.width * scale)
    if left is None:
        pic.left = int((Inches(W) - pic.width) / 2)
    else:
        pic.left = Inches(left)
    return pic


def notes(slide, s):
    slide.notes_slide.notes_text_frame.text = s.strip()


def bullets(slide, l, t, w, items, size=16, gap=11, color=INK):
    runs = []
    for it in items:
        if isinstance(it, tuple):
            s, col, bold = it
        else:
            s, col, bold = it, color, False
        runs.append((s, size, bold, col, gap))
    return text(slide, l, t, w, H - t - 0.7, runs)


# ---------------------------------------------------------------- slides -----

def build(figs, out):
    prs = new_deck()
    n = 0

    # ---------------------------------------------------------- 1 title -----
    s = blank(prs)
    rect(s, 0, 0, W, 2.45, NAVY)
    text(s, 0.9, 0.62, 11.5, 1.6,
         [("Linking Reports to Information Needs", 38, True, WHITE, 6),
          ("Pairwise attention over reports and need justifications, trained on "
           "author citations", 17, False, BLUE_ON_DARK, 0)])
    text(s, 0.9, 3.1, 11.5, 2.6,
         [("The method in one sentence", 13, True, BLUE, 8),
          ("Score every (report, need justification) pair with a "
           "DistilBERT-MNLI cross-encoder,", 20, False, INK, 2),
          ("train it with binary cross-entropy on cited positives and sampled negatives,",
           20, False, INK, 2),
          ("and serve it behind a retriever that narrows 1,000 needs to fifty.",
           20, False, INK, 0)])
    text(s, 0.9, 6.3, 11.5, 0.5,
         [("Johnny Morgan   ·   University of Maryland, Baltimore County",
           12, False, MUTED, 0)])
    notes(s, """
Opening frame. Two datasets, one function to learn. The whole talk turns on a
single design decision: we put the label on the PAIR, not on the text. Everything
else follows from that.
""")

    # ---------------------------------------------------------- 2 problem ---
    n += 1
    s = blank(prs)
    header(s, "What we have, and what we want", accent=BLUE, kicker="the problem")
    for i, (t_, b_, m_, c_) in enumerate([
        ("Reports", "r₁ … rₙ", "N ≈ 20,000", BLUE),
        ("Needs", "n₁ … nₘ, each with a written\njustification", "M > 1,000", AMBER),
        ("Citations", "an author cited report rᵢ as\nsatisfying need nⱼ",
         "≈ 30,000 edges", GREEN)]):
        l = 0.7 + i * 4.05
        rect(s, l, 2.15, 3.75, 2.15, SLATE)
        rect(s, l, 2.15, 3.75, 0.09, c_)
        text(s, l + 0.28, 2.4, 3.2, 1.8,
             [(t_, 17, True, c_, 4), (b_, 13, False, INK, 4),
              (m_, 13, True, MUTED, 0)])
    rect(s, 0.7, 4.75, 11.9, 0.75, VIOLET_S, VIOLET)
    text(s, 1.0, 4.87, 11.3, 0.6,
         [("20,000 × 1,000 = 20,000,000 candidate pairs. About 0.15% are cited.",
           15, True, INK, 0)])
    rect(s, 0.7, 5.75, 11.9, 1.1, GREEN_S, GREEN)
    text(s, 1.0, 5.88, 11.3, 0.95,
         [("The goal, both halves", 13, True, GREEN, 3),
          ("Given a new report, return the needs it satisfies. AND surface real "
           "report-need pairs nobody has cited.", 17, True, INK, 0)])
    footer(s, n, "The problem")
    notes(s, """
Say the shapes out loud. 20,000 reports, 1,000+ needs, 20 million candidate
pairs of which roughly 0.15% carry a citation.

The second half of the goal is the one that drives the design. If we only wanted
to reproduce existing citations this would be a straightforward classifier. We
want the pairs that are NOT in the citation record, and that changes what
counts as an error.
""")

    # ---------------------------------------------------------- 3 naive -----
    n += 1
    s = blank(prs)
    header(s, "The obvious approach, and why we are not using it",
           accent=BLUE, kicker="the problem")
    rect(s, 0.7, 2.1, 5.6, 3.9, SLATE)
    text(s, 1.0, 2.35, 5.0, 3.5,
         [("Treat need IDs as class labels", 17, True, MUTED, 8),
          ("Encode the report alone.", 14, False, INK, 6),
          ("Softmax over 1,000+ output classes.", 14, False, INK, 6),
          ("Done in an afternoon.", 14, False, INK, 12),
          ("But:", 14, True, RED, 6),
          ("The justifications are never read.", 14, False, RED, 5),
          ("A new need has no output unit, no training data, and no score "
           "until you retrain everything.", 14, False, RED, 5),
          ("Nothing learned transfers to an unseen need.", 14, False, RED, 0)])
    rect(s, 7.0, 2.1, 5.6, 3.9, GREEN_S, GREEN)
    text(s, 7.3, 2.35, 5.0, 3.5,
         [("Ask about the pair instead", 17, True, GREEN, 8),
          ("Feed the report AND a candidate justification into one model.",
           14, False, INK, 6),
          ("Ask a yes-or-no question: does this report satisfy this need?",
           14, False, INK, 12),
          ("Because:", 14, True, GREEN, 6),
          ("The justification is half the input.", 14, False, GREEN, 5),
          ("A new need works the moment someone writes its justification.",
           14, False, GREEN, 5),
          ("Multi-label falls out for free.", 14, False, GREEN, 0)])
    footer(s, n, "The problem")
    notes(s, """
This is the fork in the road. The left column is what most people build and it
is not wrong, it is just brittle in the one way that matters here: the need
repository changes over time and a fixed output layer cannot follow it.
""")

    # ---------------------------------------------------------- 4 divider ---
    n += 1
    s = blank(prs)
    rect(s, 0, 0, W, H, GREEN)
    text(s, 1.1, 2.9, 11.0, 1.8,
         [("Part 1", 15, True, GREEN_ON_DARK, 8),
          ("Supervised learning", 44, True, WHITE, 8),
          ("Fitting a function to labelled examples", 19, False,
           GREEN_ON_DARK, 0)])
    notes(s, "Ten minutes on fundamentals. Skip if the room already has this.")

    # ---------------------------------------------------------- 5 loop ------
    n += 1
    s = blank(prs)
    header(s, "One training step", "The loop repeats over mini-batches until "
           "held-out performance stops improving",
           accent=GREEN, kicker="supervised learning")
    figure(s, figs / "fig01_supervised_loop.png", 2.15, max_w=11.2, max_h=4.5)
    footer(s, n, "Supervised learning")
    notes(s, """
Walk it left to right. Data in, prediction out, loss measures how wrong, gradient
says which direction fixes it, update moves the weights a little.

The point to land: none of this is specific to transformers. Swap the blue box
for logistic regression and the diagram is unchanged. That is why we can talk
about the model separately from the training procedure.
""")

    # ---------------------------------------------------------- 6 split -----
    n += 1
    s = blank(prs)
    header(s, "Why we hold data back", accent=GREEN,
           kicker="supervised learning")
    cols = [("Training", "The model sees these and learns from them.", GREEN),
            ("Validation", "Never trained on. Used to choose hyperparameters "
             "and to decide when to stop.", AMBER),
            ("Test", "Touched once, at the very end, to report a number.", RED)]
    for i, (t_, d_, c_) in enumerate(cols):
        l = 0.7 + i * 4.05
        rect(s, l, 2.2, 3.75, 1.9, SLATE)
        rect(s, l, 2.2, 3.75, 0.09, c_)
        text(s, l + 0.28, 2.45, 3.2, 1.5,
             [(t_, 17, True, c_, 5), (d_, 13, False, INK, 0)])
    rect(s, 0.7, 4.55, 11.9, 1.85, RED_S, RED)
    text(s, 1.0, 4.8, 11.3, 1.5,
         [("The trap specific to this task", 15, True, RED, 5),
          ("Split by TEXT, not by pair.", 20, True, INK, 4),
          ("If a text appears in training paired with need A and in validation "
           "paired with need B, the model has already read that text. "
           "The validation score is inflated and you will not see it happen.",
           14, False, INK, 0)])
    footer(s, n, "Supervised learning")
    notes(s, """
The split-by-text point is the one people get wrong and it is silent when they
do. Nothing crashes. The number just comes out too good and the model
disappoints in production.
""")

    # ---------------------------------------------------------- 7 divider ---
    n += 1
    s = blank(prs)
    rect(s, 0, 0, W, H, AMBER_DEEP)
    text(s, 1.1, 2.9, 11.0, 1.8,
         [("Part 2", 15, True, AMBER_ON_DARK, 8),
          ("Where the label lives", 44, True, WHITE, 8),
          ("The single design decision the rest of the talk depends on",
           19, False, AMBER_ON_DARK, 0)])
    notes(s, "The conceptual core. Slow down here.")

    # ---------------------------------------------------------- 8 pairs -----
    n += 1
    s = blank(prs)
    header(s, "The label is a property of the pair",
           "Not of the text. This is the reframing.",
           accent=AMBER, kicker="where the label lives")
    figure(s, figs / "fig06_pair_construction.png", 2.15, max_w=11.0, max_h=4.5)
    footer(s, n, "Where the label lives")
    notes(s, """
Instead of "which of M classes is this text?" we ask, for each candidate,
"does this text satisfy this need, yes or no?"

The bottom box is method, not implementation detail. The sampling ratio is a
modelling choice with consequences, and it is the first thing to tune.
""")

    # ---------------------------------------------------------- 9 example ---
    n += 1
    s = blank(prs)
    header(s, "One text, three candidate needs", accent=AMBER,
           kicker="where the label lives")
    figure(s, figs / "fig10_worked_example.png", 2.05, max_w=11.2, max_h=4.6)
    footer(s, n, "Where the label lives")
    notes(s, """
Two fire, one does not. No softmax makes them compete, because each decision is
independent. That is what multi-label means mechanically.

IN-233 is the interesting row: it shares vocabulary (coastal, cities) but is
about port throughput. That is precisely the kind of distractor hard negative
mining exists to surface.
""")

    # ------------------------------------------------------ 9b citations ----
    n += 1
    s = blank(prs)
    header(s, "Our labels are citations, and that cuts both ways",
           accent=AMBER, kicker="where the labels come from")
    figure(s, figs / "fig11_citation_label_structure.png", 2.0,
           max_w=11.0, max_h=3.9)
    rect(s, 0.9, 6.05, 11.5, 0.8, RED_S, RED)
    text(s, 1.2, 6.18, 11.0, 0.65,
         [("A citation is a reliable POSITIVE. The absence of a citation is NOT "
           "a negative.", 16, True, INK, 0)])
    footer(s, n, "Where the labels come from")
    notes(s, """
The single most important slide about the data.

A citation was made by someone doing real work with real stakes, so it is a
better positive than most annotation projects produce.

But an uncited pair is three things at once: irrelevant, relevant-but-another-
report-was-cited, or relevant-and-nobody-noticed. Nothing in the data separates
them. That is positive-unlabeled learning.

If asked "so how do we train without negatives?" the answer is the next slides:
we sample assumed negatives, we are careful where we sample them, and we buy
real negatives back through adjudication.
""")

    # -------------------------------------------------- 9c PU consequences --
    n += 1
    s = blank(prs)
    header(s, "Three consequences, and each one shows up later",
           accent=AMBER, kicker="where the labels come from")
    items = [
        ("Sampled negatives are assumptions, not facts",
         "At 0.15% density the assumption is nearly always right. It is wrong most often "
         "on the pairs that look most relevant, which are exactly the ones we draw as "
         "hard negatives.", RED),
        ("Measured precision is a floor, not an estimate",
         "Flag a genuinely relevant uncited pair and the scorer calls it a false positive. "
         "True precision is at least the measured number. Model-to-model comparison is "
         "still valid; the absolute value understates.", AMBER),
        ("The errors are the product",
         "A high-scoring uncited pair is the thing you asked for. That is a workflow, "
         "not a metric problem.", GREEN),
    ]
    top = 2.2
    for t_, d_, c_ in items:
        rect(s, 0.7, top, 0.09, 1.35, c_)
        text(s, 1.05, top - 0.05, 11.4, 1.4,
             [(t_, 16.5, True, INK, 3), (d_, 13.5, False, MUTED, 0)])
        top += 1.6
    footer(s, n, "Where the labels come from")
    notes(s, """
Walk these slowly. Each one has a section in the paper.

The middle one matters for how you report results to a sponsor. If someone asks
why precision is 0.62, part of the answer is that some of those "false
positives" are correct and the citation record has not caught up.
""")

    # ------------------------------------------------------- 9d schema ------
    n += 1
    s = blank(prs)
    header(s, "How to structure the training data",
           "Three tables you curate, one the trainer derives",
           accent=AMBER, kicker="data design")
    figure(s, figs / "fig12_data_structure.png", 2.05, max_w=10.8, max_h=4.15)
    rect(s, 0.9, 6.3, 11.5, 0.55, SLATE)
    text(s, 1.2, 6.38, 11.0, 0.45,
         [("Never hand-edit the pairs table. If you do, the sampling policy stops "
           "being reproducible.", 13.5, False, INK, 0)])
    footer(s, n, "Data design")
    notes(s, """
Four fields earn their place:

cited_on   -> lets you split by TIME. You will score future reports, so hold out
              recent months and compare against the random split.
cited_by   -> lets you check the model learned a NEED and not an AUTHOR.
opened_on  -> lets you hold out needs opened after a date and test the whole
              premise: can a justification alone classify against a need the
              model never saw? If that fails, the cross-encoder is not worth
              its cost.
confidence -> optional positive weighting. Start without it.
""")

    # ---------------------------------------------------------- 10 negs -----
    n += 1
    s = blank(prs)
    header(s, "Negative sampling", "With 1,000+ needs at 0.15% density, the "
           "natural balance is about 1 positive to 660 unlabeled",
           accent=AMBER, kicker="building the training set")
    rect(s, 0.7, 2.35, 11.9, 0.95, RED_S, RED)
    text(s, 1.0, 2.52, 11.3, 0.7,
         [("Train on that directly and the model answers \"no\" to everything. "
           "99.85% accurate. Useless.", 17, True, RED, 0)])
    for i, (t_, d_, c_, bg) in enumerate([
        ("4 hard negatives", "Needs a cheap retriever ranks highly but that are "
         "NOT gold. These sit on the boundary and do the teaching.",
         RED, RED_S),
        ("4 random negatives", "Uniformly sampled. These stop the model calling "
         "everything vaguely on-topic relevant.", MUTED, SLATE)]):
        l = 0.7 + i * 6.15
        rect(s, l, 3.65, 5.75, 1.75, bg)
        text(s, l + 0.3, 3.88, 5.15, 1.4,
             [(t_ + "   per positive", 16, True, c_, 5),
              (d_, 13.5, False, INK, 0)])
    rect(s, 0.7, 5.75, 11.9, 1.15, RED_S, RED)
    text(s, 1.0, 5.88, 11.3, 0.95,
         [("The tension, stated plainly", 13, True, RED, 3),
          ("A hard negative is a pair the retriever thinks looks relevant. That is "
           "also the description of an uncited-but-relevant pair. Sharper mining "
           "means more false negatives.", 13.5, False, INK, 0)])
    footer(s, n, "Building the training set")
    notes(s, """
4:4 is a starting point, not a result.

The bottom box is the PU problem biting. Three mitigations, cheapest first:

1. Draw hard negatives from ranks 10-50, not 1-10. Genuine uncited relevance
   concentrates at the very top. Costs nothing. Start here.
2. Train hard negatives toward 0.1 instead of 0.0. Label smoothing on one class.
   The model still ranks them below positives without being told they are
   impossible.
3. Feed adjudication back. A human-confirmed irrelevant pair is a negative you
   actually know.

Whichever you pick, WRITE IT DOWN. The negative sampling policy is the most
consequential decision in the design and it is invisible in the trained weights.
""")

    # ---------------------------------------------------------- 11 divider --
    n += 1
    s = blank(prs)
    rect(s, 0, 0, W, H, VIOLET_DEEP)
    text(s, 1.1, 2.9, 11.0, 1.8,
         [("Part 3", 15, True, VIOLET_ON_DARK, 8),
          ("Transformers and attention", 44, True, WHITE, 8),
          ("How a token reads the rest of the sequence", 19, False,
           VIOLET_ON_DARK, 0)])
    notes(s, "Mechanism section. Build to the attention matrix, which is the payoff.")

    # ---------------------------------------------------------- 12 attn -----
    n += 1
    s = blank(prs)
    header(s, "Self-attention", "Every token emits a Query, a Key and a Value. "
           "Queries match against Keys to weight the Values.",
           accent=VIOLET, kicker="transformers and attention")
    figure(s, figs / "fig02_attention_mechanism.png", 2.2, max_w=11.2, max_h=4.4)
    footer(s, n, "Transformers and attention")
    notes(s, """
The soft-lookup intuition: the query says what this token is looking for, the key
says what this token offers, the value is what it contributes if selected.

The object to hold onto is Q K-transpose. It is n by n. Entry (i,j) is how much
token i attends to token j. The next slide is entirely about the shape of that
matrix.
""")

    # ---------------------------------------------------------- 13 divider --
    n += 1
    s = blank(prs)
    rect(s, 0, 0, W, H, RED_DEEP)
    text(s, 1.1, 2.9, 11.0, 1.8,
         [("Part 4", 15, True, RED_ON_DARK, 8),
          ("Pairwise attention", 44, True, WHITE, 8),
          ("What it actually means, and why it earns its cost", 19, False,
           RED_ON_DARK, 0)])
    notes(s, "The centre of the talk.")

    # ---------------------------------------------------------- 14 input ----
    n += 1
    s = blank(prs)
    header(s, "Build one sequence out of two inputs", accent=RED,
           kicker="pairwise attention")
    rect(s, 0.7, 2.3, 11.9, 1.15, NAVY)
    text(s, 0.95, 2.55, 11.4, 0.8,
         [("[CLS]  <report text>  [SEP]  <need justification>  [SEP]",
           17, True, WHITE, 0)])
    items = [
        ("[CLS]", "A special token. Its final vector is the summary of the "
         "whole pair, and it is what the classifier head reads.", BLUE),
        ("[SEP]", "Marks the boundary between the two segments.", AMBER),
        ("Segment embeddings", "Added to every token, recording which side it "
         "came from, so the model knows which half is which.", VIOLET),
    ]
    for i, (t_, d_, c_) in enumerate(items):
        top = 3.85 + i * 1.0
        rect(s, 0.7, top, 0.08, 0.8, c_)
        text(s, 1.0, top - 0.03, 11.3, 0.9,
             [(t_, 15, True, c_, 2), (d_, 13.5, False, INK, 0)])
    rect(s, 0.7, 6.85, 11.9, 0.0, SLATE)
    footer(s, n, "Pairwise attention")
    notes(s, """
The whole trick is that this is ONE sequence. Attention does not know or care
that a boundary exists; it just attends over everything. The [SEP] and segment
embeddings are how the model recovers the distinction it needs.
""")

    # ---------------------------------------------------------- 15 matrix ---
    n += 1
    s = blank(prs)
    header(s, "This is pairwise attention",
           "The two dashed blocks are the whole argument for the architecture",
           accent=RED, kicker="pairwise attention")
    figure(s, figs / "fig03_cross_segment_attention.png", 1.95,
           max_w=6.1, max_h=4.85, left=0.75)
    text(s, 7.35, 2.35, 5.3, 4.4,
         [("Four regions", 16, True, INK, 8),
          ("text reads text", 14, True, BLUE, 2),
          ("standard self-attention", 12.5, False, MUTED, 8),
          ("need reads need", 14, True, AMBER, 2),
          ("standard self-attention", 12.5, False, MUTED, 8),
          ("text ↔ justification", 14, True, RED, 2),
          ("the cross-segment blocks. \"cities\" attends to \"municipal\". "
           "\"raised\" attends to \"rise\".", 12.5, False, INK, 10),
          ("At every one of the six layers.", 14, True, RED, 8),
          ("A bi-encoder cannot do this. It has already compressed each side "
           "to one 768-d vector before any comparison happens.",
           12.5, False, MUTED, 0)])
    footer(s, n, "Pairwise attention")
    notes(s, """
THE slide. Give it time.

Point at the dashed blocks. Everything above and left of the divider is one
input, everything below and right is the other, and the off-diagonal blocks are
tokens from one reading tokens from the other.

If someone asks whether the attention weights are real: these are illustrative,
drawn to show the structure. Real weights from a trained model look messier but
have the same block structure.
""")

    # ---------------------------------------------------------- 16 compare --
    n += 1
    s = blank(prs)
    header(s, "The cost of doing it this way", accent=RED,
           kicker="pairwise attention")
    figure(s, figs / "fig04_encoder_comparison.png", 2.0, max_w=10.4, max_h=3.9)
    rect(s, 0.9, 6.1, 11.5, 0.72, GREEN_S, GREEN)
    text(s, 1.2, 6.22, 11.0, 0.6,
         [("We are not choosing. We train the cross-encoder and use a "
           "bi-encoder to decide which pairs are worth scoring.",
           14.5, True, INK, 0)])
    footer(s, n, "Pairwise attention")
    notes(s, """
The structural reason the cross-encoder cannot cache: the justification's
representation depends on which text it is paired with. There is nothing to
precompute. Every pair is a fresh forward pass.

That is the entire cost story, and it is what forces the serving design later.
""")

    # ---------------------------------------------------------- 17 model ----
    n += 1
    s = blank(prs)
    header(s, "DistilBERT-MNLI", "Three rounds of training happened before ours",
           accent=VIOLET, kicker="the model")
    figure(s, figs / "fig05_distilbert_mnli.png", 2.1, max_w=11.0, max_h=3.5)
    rect(s, 0.9, 5.85, 11.5, 1.0, VIOLET_S, VIOLET)
    text(s, 1.2, 5.97, 11.0, 0.85,
         [("Why MNLI transfers", 14, True, VIOLET, 3),
          ("MNLI asks whether a premise supports a hypothesis. We ask whether "
           "a text satisfies an information need. Same shape of question, "
           "different vocabulary.", 14, False, INK, 0)])
    footer(s, n, "The model")
    notes(s, """
Expand the acronyms as you go.

BERT: Bidirectional Encoder Representations from Transformers. 12 layers, 110M
parameters, pretrained by predicting masked tokens.

DistilBERT: the same thing compressed by knowledge distillation. 6 layers, 66M
parameters, about 60% faster. Speed matters because we run many passes per
document.

MNLI: Multi-Genre Natural Language Inference. 393,000 sentence pairs labelled
entail / contradict / neutral.

The payoff of starting from MNLI is faster convergence and better low-data
behaviour, which is exactly where the rare needs live.
""")

    # ---------------------------------------------------------- 18 divider --
    n += 1
    s = blank(prs)
    rect(s, 0, 0, W, H, GREEN_DEEP)
    text(s, 1.1, 2.9, 11.0, 1.8,
         [("Part 5", 15, True, GREEN_ON_DARK, 8),
          ("Training, evaluation, serving", 44, True, WHITE, 8),
          ("The parts that decide whether it works in practice", 19, False,
           GREEN_ON_DARK, 0)])
    notes(s, "The practical half.")

    # ---------------------------------------------------------- 19 pipeline -
    n += 1
    s = blank(prs)
    header(s, "The training pipeline", "One mini-batch, raw strings to weight update",
           accent=GREEN, kicker="training")
    figure(s, figs / "fig07_training_pipeline.png", 2.2, max_w=11.6, max_h=4.4)
    footer(s, n, "Training")
    notes(s, """
Step 2 hides a decision worth naming: truncation policy. We cut the TEXT and
always keep the justification whole, because the justification is short, curated
and reused across every pair, while the text is long and its relevant passage
could be anywhere.

If texts routinely blow the budget, window them and take the max score across
windows. Do not just raise max_length; attention cost grows with the square of
sequence length.
""")

    # ---------------------------------------------------------- 20 hparams --
    n += 1
    s = blank(prs)
    header(s, "Settings", "Starting points. The learning rate and the sampling "
           "ratio are what to tune first.", accent=GREEN, kicker="training")
    rows = [
        ("Base checkpoint", "typeform/distilbert-base-uncased-mnli"),
        ("Max sequence length", "256 tokens, truncate the text only"),
        ("Batch size", "32 pairs"),
        ("Optimiser", "AdamW, weight decay 0.01"),
        ("Learning rate", "2e-5, linear decay, 10% warmup"),
        ("Epochs", "3, early stopping on validation mAP"),
        ("Loss", "binary cross-entropy on the sigmoid output"),
        ("Negative sampling", "4 hard + 4 random per positive"),
        ("Precision", "fp16"),
    ]
    tbl = s.shapes.add_table(len(rows) + 1, 2, Inches(0.9), Inches(2.15),
                             Inches(11.5), Inches(0.4)).table
    tbl.columns[0].width = Inches(3.6)
    tbl.columns[1].width = Inches(7.9)
    for j, h_ in enumerate(("Setting", "Value")):
        c = tbl.cell(0, j)
        c.text = h_
        p = c.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.name = FONT
        p.font.color.rgb = WHITE
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
    for i, (k_, v_) in enumerate(rows, start=1):
        for j, val in enumerate((k_, v_)):
            c = tbl.cell(i, j)
            c.text = val
            p = c.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.name = FONT
            p.font.bold = (j == 0)
            p.font.color.rgb = INK
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 else SLATE
    footer(s, n, "Training")
    notes(s, """
Three epochs is the usual range for fine-tuning a pretrained encoder. Past that
the model starts overwriting the pretrained representations rather than adapting
them, and you see it as validation performance falling while training loss keeps
dropping.
""")

    # ---------------------------------------------------------- 21 eval -----
    n += 1
    s = blank(prs)
    header(s, "Evaluating a multi-label model",
           "Accuracy is unusable. Predicting nothing scores above 99%.",
           accent=GREEN, kicker="evaluation")
    figure(s, figs / "fig08_multilabel_eval.png", 2.0, max_w=10.8, max_h=3.65)
    rect(s, 0.9, 5.9, 11.5, 0.95, GREEN_S, GREEN)
    text(s, 1.2, 6.02, 11.0, 0.8,
         [("The single most informative number is the GAP between micro-F1 and "
           "macro-F1.", 15, True, INK, 2),
          ("A large gap means the model works on common needs and fails on rare "
           "ones. Micro hides the long tail. Macro exposes it.",
           13.5, False, INK, 0)])
    footer(s, n, "Evaluation")
    notes(s, """
Micro-F1 pools every pair decision, so frequent needs dominate. Macro-F1 averages
per-need F1, so a need with six examples counts as much as one with nine hundred.

Use mAP for model selection because it evaluates the ranking and does not depend
on a threshold you have not tuned yet.

The threshold is a product decision. Human reviews everything, favour recall.
Acted on automatically, favour precision. Tune it on validation, never test.
""")

    # ---------------------------------------------------------- 22 serving --
    n += 1
    s = blank(prs)
    header(s, "Serving against 1,000+ needs",
           "1,000 needs x 5 ms per pair = 5 s per report, so ~28 GPU-hours to "
           "sweep the 20,000-report backlog once",
           accent=BLUE, kicker="serving")
    figure(s, figs / "fig09_serving.png", 2.1, max_w=11.0, max_h=3.55)
    rect(s, 0.9, 5.9, 11.5, 0.95, RED_S, RED)
    text(s, 1.2, 6.02, 11.0, 0.8,
         [("Stage one sets a hard ceiling on recall.", 15, True, RED, 2),
          ("If the right need is not in the top k, no amount of cross-encoder "
           "quality recovers it. Measure recall@k first and report it next to "
           "the end-to-end number.", 13.5, False, INK, 0)])
    footer(s, n, "Serving")
    notes(s, """
51 forward passes plus a vector search, and the cost does not grow as the
repository grows. That is the whole point.

The recall ceiling is the thing teams forget. If recall@50 is 0.88 then 0.88 is
the best the full system can do, and tuning stage two past that is wasted work.
""")

    # ------------------------------------------------------ 22b discovery ---
    n += 1
    s = blank(prs)
    header(s, "The discovery loop",
           "The uncited pairs the model scores highly are the output, not the error",
           accent=VIOLET, kicker="discovery")
    figure(s, figs / "fig13_discovery_loop.png", 2.1, max_w=11.0, max_h=3.7)
    rect(s, 0.9, 6.0, 11.5, 0.9, VIOLET_S, VIOLET)
    text(s, 1.2, 6.12, 11.0, 0.75,
         [("Both verdicts are worth more than a sampled label.", 14.5, True, INK, 2),
          ("YES is a discovery and a new positive edge. NO is a CONFIRMED negative, "
           "which the raw data never contained.", 13.5, False, INK, 0)])
    footer(s, n, "Discovery")
    notes(s, """
This is the answer to "the metric punishes us for succeeding."

Stop treating the uncited region as noise and start treating it as a queue.
Score it, rank it, send the top few hundred to a human. Every verdict is a label
that did not exist before, and adjudication is the ONLY mechanism that shrinks
the unknown region.

One caution: keep a fixed evaluation set of citations that never receives
model-proposed edges, or the model starts grading its own homework.
""")

    # ------------------------------------------------- 22c discovery metrics
    n += 1
    s = blank(prs)
    header(s, "Measuring discovery, since standard metrics will not",
           "Precision and recall answer \"does it reproduce what authors already "
           "wrote?\" That is not the question you asked.",
           accent=VIOLET, kicker="discovery")
    items = [
        ("Adjudicated precision @ k",
         "Of the top k uncited pairs sent to a human, what fraction were confirmed real? "
         "The direct measure of discovery yield, and the number to report to a sponsor. "
         "Costs human time and nothing substitutes for it.", VIOLET),
        ("Held-out citation recovery",
         "Hide 10% of known citations, train without them, check where they rank among "
         "that report's uncited pairs. A proxy for discovery that costs NO human time, "
         "so it is what you run first.", GREEN),
        ("Yield curve",
         "Confirmed discoveries against pairs adjudicated. It will flatten. Where it "
         "flattens tells you when to stop reviewing.", BLUE),
    ]
    top = 2.35
    for t_, d_, c_ in items:
        rect(s, 0.7, top, 0.09, 1.3, c_)
        text(s, 1.05, top - 0.05, 11.4, 1.35,
             [(t_, 16.5, True, INK, 3), (d_, 13.5, False, MUTED, 0)])
        top += 1.55
    footer(s, n, "Discovery")
    notes(s, """
Lead with the middle one operationally. Held-out citation recovery needs no
adjudication budget, so it runs first, right after training.

Adjudicated precision @ k is what you report upward. It is also the only number
that directly answers "did this find anything new."
""")

    # ---------------------------------------------------------- 23 risks ----
    n += 1
    s = blank(prs)
    header(s, "What we expect to go wrong", accent=RED, kicker="limitations")
    risks = [
        ("Justification quality is the ceiling",
         "A need described in one vague sentence will perform badly no matter "
         "how long we train. Sort needs by F1 and read the worst justifications "
         "before blaming the model."),
        ("Citation is a proxy for relevance, and proxies drift",
         "An author cited a report because it was relevant AND because they found it, "
         "could access it, and remembered it. The model learns all four. Check "
         "per-need precision against report age and indexing tier."),
        ("Circularity",
         "If a justification was written by reading the reports now cited against it, "
         "the model scores well for a reason that will not survive new documents. "
         "Compare needs.opened_on against citations.cited_on."),
        ("Preferential attachment",
         "A few needs attract most citations. The model learns that skew and "
         "reproduces it, making popular needs look well-served regardless of the "
         "texts. Plot per-need F1 against citation count."),
    ]
    top = 2.15
    for t_, d_ in risks:
        rect(s, 0.7, top, 0.08, 1.05, RED)
        text(s, 1.0, top - 0.05, 11.4, 1.1,
             [(t_, 15.5, True, INK, 2), (d_, 13, False, MUTED, 0)])
        top += 1.17
    footer(s, n, "Limitations")
    notes(s, """
Circularity is the one to worry about most, because it is the failure mode that
produces a GOOD number and a BAD system. Everything else announces itself
eventually.
""")

    # ---------------------------------------------------------- 24 close ----
    n += 1
    s = blank(prs)
    rect(s, 0, 0, W, 2.3, NAVY)
    text(s, 0.9, 0.72, 11.5, 1.3,
         [("Suggested order of work", 32, True, WHITE, 0)])
    steps = [
        ("1", "Build the three tables. Count the citations and measure the TRUE "
              "density first, since every estimate keys off it.", BLUE),
        ("2", "Run the frozen bi-encoder alone and measure recall@50. One afternoon, "
              "no training, and it tells you the ceiling.", AMBER),
        ("3", "Train the cross-encoder. A few hours on one GPU. This is the cheap part.", GREEN),
        ("4", "Run held-out citation recovery. Costs no human time.", GREEN),
        ("5", "Only then adjudicate a few hundred top-ranked uncited pairs.", VIOLET),
    ]
    top = 2.55
    for num, d_, c_ in steps:
        rect(s, 0.9, top, 0.45, 0.62, c_)
        text(s, 0.9, top + 0.06, 0.45, 0.5, [(num, 16, True, WHITE, 0)],
             align=PP_ALIGN.CENTER)
        text(s, 1.55, top + 0.04, 10.9, 0.6, [(d_, 14.5, False, INK, 0)])
        top += 0.78
    rect(s, 0.9, 6.5, 11.5, 0.42, RED_S, RED)
    text(s, 1.2, 6.55, 11.0, 0.35,
         [("Report the micro/macro gap, recall@k, and adjudicated precision@k. "
           "Not the headline F1.", 13.5, True, INK, 0)])
    footer(s, n, "Summary")
    notes(s, """
Close on sequencing, because the order saves money.

Step 2 is the one people skip. Measuring the retrieval ceiling before training
anything costs an afternoon and can tell you the whole approach is capped lower
than you need, before you have spent a GPU-hour or an adjudication hour.

Step 5 last, because adjudication is the expensive resource in this project.
Compute is not.
""")

    prs.save(str(out))
    return len(prs.slides), out


# ---------------------------------------------------------------- main -------

def main():
    here = Path(__file__).resolve().parent
    ap = argparse.ArgumentParser()
    ap.add_argument("--figures", default=str(here / "figures"))
    ap.add_argument("--out", default=str(
        here / "Information_Need_Cross_Encoder_Briefing.pptx"))
    args = ap.parse_args()

    figs = Path(args.figures)
    missing = [f"fig{i:02d}" for i in range(1, 14)
               if not list(figs.glob(f"fig{i:02d}_*.png"))]
    if missing:
        raise SystemExit(f"missing figures: {missing}. Run make_figures.py first.")

    out = Path(args.out)
    out.parent.mkdir(parents=True, exist_ok=True)
    n, path = build(figs, out)
    print(f"{n} slides -> {path}  ({path.stat().st_size / 1024:.0f} KB)")


if __name__ == "__main__":
    main()
