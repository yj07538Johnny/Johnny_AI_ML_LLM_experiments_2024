#!/usr/bin/env python3
"""deckkit — 16:9 slide primitives for python-pptx, sharing figkit's palette.

The pptx companion to figkit. figkit draws the figures; deckkit lays out the
slides they sit on, using the same hues so a colour means one thing in both.

    from deckkit import (Deck, C, INK, BLUE, GREEN_S,
                         rect, text, header, footer, figure, notes, table)

    d = Deck()
    s = d.slide()
    header(s, "Headline", "subtitle", accent=BLUE, kicker="section")
    figure(s, "figures/fig2.png", top=2.1, max_h=4.0)
    footer(s, 1, "Section")
    notes(s, "what to say out loud")
    d.save("out.pptx")

Coordinates are inches on a 13.333 x 7.5 canvas. Keep content boxes ending
above y=6.9 or they collide with the footer, a mistake that is invisible in
source and obvious the moment you render the deck to PDF and look at it.
"""

from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation                                # noqa: E402
from pptx.dml.color import RGBColor                          # noqa: E402
from pptx.enum.shapes import MSO_SHAPE                       # noqa: E402
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN              # noqa: E402
from pptx.util import Inches, Pt                             # noqa: E402

import figkit                                                # noqa: E402

W, H = 13.333, 7.5          # inches, 16:9
FONT = "Calibri"
FOOTER_Y = H - 0.52
SAFE_BOTTOM = 6.9           # content must end above this


def C(hexstr) -> RGBColor:
    """figkit hex -> pptx RGBColor."""
    return RGBColor(*figkit.rgb(hexstr))


INK = C(figkit.INK)
MUTED = C(figkit.MUTED)
WHITE = C(figkit.PAPER)
SLATE = C(figkit.SLATE_L)
BLUE = C(figkit.BLUE)
AMBER = C(figkit.AMBER)
GREEN = C(figkit.GREEN)
RED = C(figkit.RED)
VIOLET = C(figkit.VIOLET)

NAVY = C(figkit.INK_DARK)
GREEN_S = C(figkit.GREEN_S)
RED_S = C(figkit.RED_S)
VIOLET_S = C(figkit.VIOLET_S)

BLUE_ON_DARK = C(figkit.BLUE_ON_DARK)
GREEN_ON_DARK = C(figkit.GREEN_ON_DARK)
AMBER_ON_DARK = C(figkit.AMBER_ON_DARK)
VIOLET_ON_DARK = C(figkit.VIOLET_ON_DARK)
RED_ON_DARK = C(figkit.RED_ON_DARK)


class Deck:
    """A 16:9 presentation with blank layouts."""

    def __init__(self, width=W, height=H):
        self.prs = Presentation()
        self.prs.slide_width = Inches(width)
        self.prs.slide_height = Inches(height)

    def slide(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def save(self, path):
        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))
        return len(self.prs.slides), path


def rect(slide, l, t, w, h, color, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t),
                                Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """runs = [(string, size_pt, bold, RGBColor, space_after_pt), ...]"""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (s, sz, bold, col, after) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = s
        p.alignment = align
        p.space_after = Pt(after)
        f = p.font
        f.name = FONT
        f.size = Pt(sz)
        f.bold = bold
        f.color.rgb = col
    return tb


def header(slide, title, sub=None, accent=BLUE, kicker=None):
    """Left accent bar, optional kicker, headline, optional subtitle."""
    rect(slide, 0, 0, 0.22, H, accent)
    y = 0.42
    if kicker:
        text(slide, 0.62, y, 11.8, 0.34, [(kicker.upper(), 11, True, accent, 0)])
        y += 0.38
    runs = [(title, 30, True, INK, 4)]
    if sub:
        runs.append((sub, 15, False, MUTED, 0))
    text(slide, 0.62, y, 11.9, 1.25, runs)


def footer(slide, n, label):
    text(slide, 0.62, FOOTER_Y, 9.0, 0.3, [(label, 10, False, MUTED, 0)])
    text(slide, 11.4, FOOTER_Y, 1.3, 0.3, [(str(n), 10, False, MUTED, 0)],
         align=PP_ALIGN.RIGHT)


def figure(slide, path, top, max_w=11.4, max_h=4.5, left=None):
    """Insert a PNG scaled to fit (max_w, max_h). Centred unless `left` given."""
    pic = slide.shapes.add_picture(str(path), Inches(0), Inches(top),
                                   width=Inches(max_w))
    if pic.height > Inches(max_h):
        scale = Inches(max_h) / pic.height
        pic.height = Inches(max_h)
        pic.width = int(pic.width * scale)
    pic.left = int((Inches(W) - pic.width) / 2) if left is None else Inches(left)
    return pic


def notes(slide, s):
    slide.notes_slide.notes_text_frame.text = s.strip()


def table(slide, l, t, w, rows, col_widths=None, header_fill=NAVY,
          fs=12, header_fs=12.5, row_h=0.4):
    """rows[0] is the header. Returns the pptx table.

    col_widths are inches and must sum to w; omit for equal columns.
    """
    n_rows, n_cols = len(rows), len(rows[0])
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(l), Inches(t),
                                   Inches(w), Inches(row_h))
    tbl = shape.table
    if col_widths:
        for j, cw in enumerate(col_widths):
            tbl.columns[j].width = Inches(cw)

    for j, val in enumerate(rows[0]):
        c = tbl.cell(0, j)
        c.text = str(val)
        p = c.text_frame.paragraphs[0]
        p.font.bold, p.font.size, p.font.name = True, Pt(header_fs), FONT
        p.font.color.rgb = WHITE
        c.fill.solid()
        c.fill.fore_color.rgb = header_fill

    for i, row in enumerate(rows[1:], start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = str(val)
            p = c.text_frame.paragraphs[0]
            p.font.size, p.font.name = Pt(fs), FONT
            p.font.bold = (j == 0)
            p.font.color.rgb = INK
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 else SLATE
    return tbl
